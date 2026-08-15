import io
import re
from pathlib import Path

from .helpers import get_output_dir, resolve_output_path

PIL_FORMAT = {"JPG": "JPEG", "PNG": "PNG", "WEBP": "WEBP"}


def _page_count(src):
    import pymupdf as fitz

    with fitz.open(src) as doc:
        return len(doc)


def pdf_to_txt(src, out):
    import pymupdf as fitz

    with fitz.open(src) as doc:
        text = "\n\n".join(page.get_text() for page in doc)
    Path(out).write_text(text, encoding="utf-8")


def pdf_to_md(src, out):
    import pymupdf as fitz

    with fitz.open(src) as doc:
        text = "\n\n".join(page.get_text() for page in doc)
    text = re.sub(r"\n{3,}", "\n\n", text)
    Path(out).write_text(text, encoding="utf-8")


def pdf_to_html(src, out):
    import pymupdf as fitz

    with fitz.open(src) as doc:
        body = "\n".join(page.get_text("html") for page in doc)
    html = (
        "<!DOCTYPE html>\n<html><head><meta charset='utf-8'><title>PDF</title>"
        "<style>body{font-family:Georgia,serif;font-size:12pt;margin:2cm}</style></head>"
        f"<body>{body}</body></html>\n"
    )
    Path(out).write_text(html, encoding="utf-8")


def pdf_to_images(src, out_paths, target_fmt, dpi=150):
    import pymupdf as fitz
    from PIL import Image

    fmt = PIL_FORMAT[target_fmt]
    with fitz.open(src) as doc:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=dpi)
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            img.save(out_paths[i], format=fmt)


def pdf_to_docx(src, out):
    import pymupdf as fitz
    from docx import Document
    from docx.enum.text import WD_BREAK

    document = Document()
    first = True
    with fitz.open(src) as doc:
        for page in doc:
            if not first:
                document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
            first = False
            for block in page.get_text("dict").get("blocks", []):
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    text = ""
                    bold = False
                    italic = False
                    for span in line.get("spans", []):
                        fname = span.get("font", "")
                        flags = span.get("flags", 0)
                        bold = bold or "bold" in fname.lower() or bool(flags & 16)
                        italic = italic or "italic" in fname.lower() or bool(flags & 2)
                        text += span.get("text", "")
                    if text.strip():
                        run = document.add_paragraph().add_run(text)
                        run.bold = bold
                        run.italic = italic
    document.save(out)


def pdf_to_pptx(src, out, dpi=120):
    import pymupdf as fitz
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    with fitz.open(src) as doc:
        if doc:
            page0 = doc[0]
            ratio = page0.rect.width / page0.rect.height
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(10 / ratio)
        for page in doc:
            pix = page.get_pixmap(dpi=dpi)
            buf = io.BytesIO(pix.tobytes("png"))
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out)


def split_pdf_every(src, pages_per_file):
    """Split a PDF into parts of pages_per_file pages each. Returns output paths."""
    import pymupdf as fitz

    out_dir = get_output_dir(src)
    base = Path(src).stem
    paths = []
    with fitz.open(src) as doc:
        total = len(doc)
        part = 1
        start = 0
        while start < total:
            end = min(start + pages_per_file, total)
            out = out_dir / f"{base}_part{part}.pdf"
            new = fitz.open()
            new.insert_pdf(doc, from_page=start, to_page=end - 1)
            new.save(out)
            new.close()
            paths.append(out)
            start = end
            part += 1
    return paths


def split_pdf_ranges(src, ranges):
    """Split a PDF by 1-based inclusive page ranges [(start, end), ...]."""
    import pymupdf as fitz

    out_dir = get_output_dir(src)
    base = Path(src).stem
    paths = []
    with fitz.open(src) as doc:
        total = len(doc)
        for i, (start, end) in enumerate(ranges, 1):
            start = max(1, min(start, total))
            end = max(start, min(end, total))
            out = out_dir / f"{base}_part{i}.pdf"
            new = fitz.open()
            new.insert_pdf(doc, from_page=start - 1, to_page=end - 1)
            new.save(out)
            new.close()
            paths.append(out)
    return paths


def merge_pdfs(src_paths, out_path):
    """Merge multiple PDFs into one, in order. Returns the output path."""
    import pymupdf as fitz

    merged = fitz.open()
    for p in src_paths:
        with fitz.open(p) as d:
            merged.insert_pdf(d)
    merged.save(out_path)
    merged.close()
    return Path(out_path)


def compress_pdf(src, out, quality=60, dpi=110):
    """Re-render every page as a JPEG and rebuild the PDF to shrink it."""
    import pymupdf as fitz
    from PIL import Image

    with fitz.open(src) as src_doc:
        new_doc = fitz.open()
        try:
            for page in src_doc:
                pix = page.get_pixmap(dpi=dpi)
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=quality, optimize=True)
                new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                new_page.insert_image(new_page.rect, stream=buf.getvalue())
            new_doc.save(out, deflate=True)
        finally:
            new_doc.close()


def _image_handler(target_fmt):
    def handler(src, out, on_conflict=None):
        n = _page_count(src)
        paths = [
            resolve_output_path(src, target_fmt, on_conflict=on_conflict, suffix=f"_page{i + 1}")
            for i in range(n)
        ]
        pdf_to_images(src, paths, target_fmt)
        return paths

    return handler


HANDLERS = {
    ("PDF", "DOCX"): pdf_to_docx,
    ("PDF", "PPTX"): pdf_to_pptx,
    ("PDF", "TXT"): pdf_to_txt,
    ("PDF", "MD"): pdf_to_md,
    ("PDF", "HTML"): pdf_to_html,
    ("PDF", "JPG"): _image_handler("JPG"),
    ("PDF", "PNG"): _image_handler("PNG"),
    ("PDF", "WEBP"): _image_handler("WEBP"),
}