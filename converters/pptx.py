import os
import tempfile
from html import escape
from pathlib import Path

from .helpers import html_to_pdf, page_html, resolve_output_path
from .pdf import _page_count, pdf_to_images


def _slides_text(src):
    from pptx import Presentation

    prs = Presentation(src)
    slides = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text:
                    parts.append(text)
            elif shape.has_table:
                rows = [
                    " | ".join(cell.text for cell in row.cells)
                    for row in shape.table.rows
                ]
                if rows:
                    parts.append("\n".join(rows))
        slides.append(parts)
    return slides


def pptx_to_txt(src, out):
    blocks = []
    for i, parts in enumerate(_slides_text(src), 1):
        if parts:
            blocks.append(f"Slide {i}\n" + "\n".join(parts))
    Path(out).write_text("\n\n".join(blocks) + "\n", encoding="utf-8")


def pptx_to_md(src, out):
    blocks = []
    for i, parts in enumerate(_slides_text(src), 1):
        if parts:
            blocks.append(f"## Slide {i}\n\n" + "\n\n".join(parts))
    Path(out).write_text("\n\n".join(blocks) + "\n", encoding="utf-8")


def _pptx_html(src):
    body = []
    for i, parts in enumerate(_slides_text(src), 1):
        if not parts:
            continue
        div = "".join(f"<p>{escape(line)}</p>" for line in parts)
        body.append(f"<div style='page-break-before:always'>{div}</div>")
    return page_html("".join(body))


def pptx_to_pdf(src, out):
    html_to_pdf(_pptx_html(src), out)


def _pptx_to_images(src, target_fmt, on_conflict=None):
    fd, tmp = tempfile.mkstemp(suffix=".pdf")
    os.close(fd)
    try:
        html_to_pdf(_pptx_html(src), tmp)
        n = _page_count(tmp)
        paths = [
            resolve_output_path(src, target_fmt, on_conflict=on_conflict, suffix=f"_page{i + 1}")
            for i in range(n)
        ]
        pdf_to_images(tmp, paths, target_fmt)
        return paths
    finally:
        Path(tmp).unlink(missing_ok=True)


def _image_handler(target_fmt):
    def handler(src, out, on_conflict=None):
        return _pptx_to_images(src, target_fmt, on_conflict=on_conflict)

    return handler


HANDLERS = {
    ("PPTX", "PDF"): pptx_to_pdf,
    ("PPTX", "TXT"): pptx_to_txt,
    ("PPTX", "MD"): pptx_to_md,
    ("PPTX", "JPG"): _image_handler("JPG"),
    ("PPTX", "PNG"): _image_handler("PNG"),
}